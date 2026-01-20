"""Document parsing service for various file formats."""

import io
import json
from typing import Dict, Any, Optional
from pathlib import Path

from services.markdown_converter import (
    get_file_category,
    convert_to_markdown,
    get_supported_formats as get_converter_formats,
    ALL_SUPPORTED as ALL_SUPPORTED_EXTENSIONS
)


def parse_text(content: bytes, filename: str) -> Dict[str, Any]:
    """Parse plain text file."""
    text = content.decode("utf-8")
    return {
        "text": text,
        "filename": filename,
        "format": "text",
        "metadata": {
            "char_count": len(text),
            "line_count": text.count("\n") + 1
        }
    }


def parse_markdown(content: bytes, filename: str) -> Dict[str, Any]:
    """Parse markdown file."""
    text = content.decode("utf-8")
    return {
        "text": text,
        "filename": filename,
        "format": "markdown",
        "metadata": {
            "char_count": len(text),
            "line_count": text.count("\n") + 1,
            "has_headers": "#" in text
        }
    }


def parse_json(content: bytes, filename: str) -> Dict[str, Any]:
    """Parse JSON file and convert to readable text."""
    data = json.loads(content.decode("utf-8"))

    if isinstance(data, dict):
        text = json.dumps(data, indent=2)
    elif isinstance(data, list):
        text = "\n\n".join(json.dumps(item, indent=2) for item in data)
    else:
        text = str(data)

    return {
        "text": text,
        "filename": filename,
        "format": "json",
        "metadata": {
            "char_count": len(text),
            "is_array": isinstance(data, list),
            "keys": list(data.keys()) if isinstance(data, dict) else None
        }
    }


def parse_pdf(content: bytes, filename: str) -> Dict[str, Any]:
    """Parse PDF file and extract text."""
    try:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(content))
        pages = []

        for i, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text:
                pages.append({
                    "page": i + 1,
                    "text": page_text
                })

        full_text = "\n\n".join(p["text"] for p in pages)

        return {
            "text": full_text,
            "filename": filename,
            "format": "pdf",
            "metadata": {
                "char_count": len(full_text),
                "page_count": len(reader.pages),
                "pages": pages
            }
        }
    except Exception as e:
        raise ValueError(f"Failed to parse PDF: {str(e)}")


def parse_docx(content: bytes, filename: str) -> Dict[str, Any]:
    """Parse DOCX file and extract text."""
    try:
        from docx import Document

        doc = Document(io.BytesIO(content))
        paragraphs = []

        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)

        full_text = "\n\n".join(paragraphs)

        return {
            "text": full_text,
            "filename": filename,
            "format": "docx",
            "metadata": {
                "char_count": len(full_text),
                "paragraph_count": len(paragraphs)
            }
        }
    except Exception as e:
        raise ValueError(f"Failed to parse DOCX: {str(e)}")


def parse_pptx(content: bytes, filename: str) -> Dict[str, Any]:
    """Parse PowerPoint file and extract text."""
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content))
        slides = []

        for i, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)

            if slide_texts:
                slides.append({
                    "slide": i + 1,
                    "text": "\n".join(slide_texts)
                })

        # Format as markdown with slide headers
        formatted_slides = []
        for slide in slides:
            formatted_slides.append(f"## Slide {slide['slide']}\n\n{slide['text']}")

        full_text = "\n\n---\n\n".join(formatted_slides)

        return {
            "text": full_text,
            "filename": filename,
            "format": "pptx",
            "metadata": {
                "char_count": len(full_text),
                "slide_count": len(prs.slides),
                "slides": slides
            }
        }
    except Exception as e:
        raise ValueError(f"Failed to parse PowerPoint: {str(e)}")


def parse_convertible(content: bytes, filename: str) -> Dict[str, Any]:
    """Parse a file by converting it to markdown first."""
    try:
        text_content = content.decode("utf-8")
    except UnicodeDecodeError:
        # Try with latin-1 as fallback
        text_content = content.decode("latin-1")

    result = convert_to_markdown(text_content, filename)
    markdown_content = result.get("content", text_content)

    return {
        "text": markdown_content,
        "filename": filename,
        "format": "converted",
        "metadata": {
            "original_format": result.get("converted_from", "unknown"),
            "language": result.get("language"),
            "char_count": len(markdown_content)
        }
    }


def parse_document(content: bytes, filename: str) -> Dict[str, Any]:
    """
    Parse a document based on its file extension.

    Args:
        content: File content as bytes
        filename: Original filename

    Returns:
        Dict with parsed text and metadata
    """
    ext = Path(filename).suffix.lower()
    name = Path(filename).name

    # Native parsers for specific formats
    native_parsers = {
        ".txt": parse_text,
        ".md": parse_markdown,
        ".markdown": parse_markdown,
        ".json": parse_json,
        ".pdf": parse_pdf,
        ".docx": parse_docx,
        ".pptx": parse_pptx,
        ".ppt": parse_pptx,  # Try pptx parser for legacy files
    }

    # Check file category
    category, lang, desc = get_file_category(filename)

    if category == 'native' and ext in native_parsers:
        parser = native_parsers[ext]
        result = parser(content, filename)
    elif category == 'convert':
        result = parse_convertible(content, filename)
    elif ext in native_parsers:
        # Fallback for native formats
        parser = native_parsers[ext]
        result = parser(content, filename)
    else:
        raise ValueError(f"Unsupported file format: {ext}. Supported: {', '.join(sorted(ALL_SUPPORTED_EXTENSIONS))}")

    # Transform to consistent format expected by frontend
    text = result.get("text", "")
    metadata = result.get("metadata", {})

    return {
        "filename": result.get("filename", filename),
        "content": text,
        "metadata": {
            "file_type": ext.lstrip("."),
            "char_count": len(text),
            "word_count": len(text.split()),
            "line_count": text.count("\n") + 1,
            "converted": category == 'convert',
            "original_format": metadata.get("original_format"),
            "language": metadata.get("language")
        }
    }


def get_supported_formats() -> Dict[str, Any]:
    """Get comprehensive list of supported file formats."""
    return get_converter_formats()
